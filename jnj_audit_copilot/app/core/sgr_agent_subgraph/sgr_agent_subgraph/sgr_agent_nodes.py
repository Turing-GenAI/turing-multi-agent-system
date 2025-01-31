import os

from langchain.prompts import PromptTemplate
from langchain_core.messages import AIMessage
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from ....common.config import graph_inputs
from ....common.constants import (
    SGR_FILEPATHS,
    SGR_OUTPUT_FOLDER,
    SGR_PPT_FILENAME,
    bold_end,
    bold_start,
    trial_id,
)
from ....common.descriptions import PD_Term_DVTERM, pd_terms
from ....prompt_hub.sgr_pd_agent_prompts import sgr_pd_agent_prompts
from ....utils.helpers import read_file
from ....utils.langchain_azure_openai import azure_chat_openai_client as llm
from ....utils.log_setup import get_logger
from ....utils.state_definitions import SGRSubGraphState
from .sgr_nodes_functions import createSGRPresentationFunctions, sgrPDNodesFunctions

# Get the same logger instance set up earlier
logger = get_logger()

if not os.path.exists(os.path.dirname(SGR_FILEPATHS["intermediate"]["dvdecode"])):
    os.makedirs(os.path.dirname(SGR_FILEPATHS["intermediate"]["dvdecode"]))

TRIAL_ID = graph_inputs["trigger_list"][0]["trial_id"]


class SGRSubGraphNodes:
    def __init__(self, tid=trial_id):
        logger.debug("Initialising SGRSubGraphNodes .... ")
        self.tid = tid
        self.create_presentation_functions = createSGRPresentationFunctions()

    def fetch_SGR_data(self, state: SGRSubGraphState):
        logger.debug("Calling function : fetch_sgr_data..")
        input_filepaths = {
            "pd_synthetic_data_filepath": SGR_FILEPATHS["input"]["pd_synthetic_data_filepath"],
            "quasar_data_filepath": SGR_FILEPATHS["input"]["quasar_data"],
            "sponsor_inspection_data_filepath": SGR_FILEPATHS["input"]["sponsor_inspections_data"],
            "sqi_data_filepath": SGR_FILEPATHS["input"]["sqi_data"],
            "serious_breach_data_filepath": SGR_FILEPATHS["input"]["serious_breach_data"],
            "qis_sbs_data_filepath": SGR_FILEPATHS["input"]["qis_sbs_data"],
        }
        return {
            "input_filepath": input_filepaths,
            "sgr_agent_messages": AIMessage(
                name=f"{bold_start}CRM - fetch_sgr_data node:{bold_end}",
                content="Initiating retrieval of necessary data for SGR " "Executive Summaries. ",
            ),
        }

    def create_sgr_presentation(self, state: SGRSubGraphState):
        # Load the two CSV files for Slide 1
        df_major_devs = read_file(SGR_FILEPATHS["intermediate"]["dvdecode"], "csv")
        df_other_devs = read_file(SGR_FILEPATHS["intermediate"]["pdterm"], "csv")
        sgr_keyobs_path = SGR_FILEPATHS["intermediate"]["keyobs"]
        try:
            # Read the GPT-4 response from the text file
            with open(SGR_FILEPATHS["intermediate"]["keyobs"], "r") as file:
                gpt4_response = file.read()
        except Exception as e:
            logger.error(f"No content found in {sgr_keyobs_path}. Due to error: {e}")
            gpt4_response = "NA"
        gpt4_response = f"Key Observations/Comments:\n\n {gpt4_response}"

        # Create PowerPoint presentation
        prs = Presentation()

        # Slide 1: Create layout and add tables for Major Protocol Deviations
        # and GPT-4 response
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title to the slide
        title_shape = slide.shapes.title
        title_shape.text = "Major Protocol Deviation Trials 12/3"
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.5)
        title_shape.width = Inches(4)
        title_shape.height = Inches(0.5)
        title_shape.text_frame.paragraphs[0].font.size = Pt(15)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)  # Dark blue color

        # Add the first DataFrame to the left side of the slide
        left_table = self.create_presentation_functions.add_table_to_slide(
            slide,
            df_major_devs,
            Inches(0.5),
            Inches(1),
            Inches(4.1),
            Inches(2),
            header_color=RGBColor(0, 0, 139),
            row_color=RGBColor(173, 216, 230),
        )
        left_table.columns[0].width = Inches(3.1)  # Custom width for the first column
        left_table.columns[1].width = Inches(1)  # Custom width for the second column

        # Add the second DataFrame to the right side of the slide
        right_table = self.create_presentation_functions.add_table_to_slide(
            slide,
            df_other_devs,
            Inches(4.7),
            Inches(1),
            Inches(5.1),
            Inches(5.6),
            header_color=RGBColor(0, 0, 139),
            row_color=RGBColor(173, 216, 230),
        )
        right_table.columns[0].width = Inches(4)  # Custom width for the first column
        right_table.columns[1].width = Inches(1)  # Custom width for the second column

        # Add the GPT-4 Response in the left bottom corner
        left_bottom_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(4.1), Inches(2.7))
        left_bottom_text_frame = left_bottom_textbox.text_frame
        left_bottom_text_frame.word_wrap = True
        left_bottom_text_frame.text = gpt4_response
        left_bottom_textbox.fill.solid()
        left_bottom_textbox.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Set light blue background
        # Set the first line bold
        first_paragraph = left_bottom_text_frame.paragraphs[0]
        first_run = first_paragraph.runs[0]
        first_run.font.bold = True
        first_run.font.size = Pt(10)  # Set font size to 10 points for the first line

        # Set the font size for the rest of the text
        for paragraph in left_bottom_text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(10)

        # Slide 2: Executive Summary
        sections = [
            (
                "Major Protocol Deviations",
                state["intermediate_files"]["pd_section"],
            ),
            (
                "Significant Issue Escalations",
                state["intermediate_files"]["significant_issue_escalations_section"],
            ),
            (
                "Site and Sponsor Inspections",
                state["intermediate_files"]["site_sponsor_inspections_section"],
            ),
            ("QA Audits", state["intermediate_files"]["qa_audit_section"]),
        ]

        # Collect all the content for each section
        sections_content = []
        for title, filename in sections:
            try:
                with open(filename, "r") as file:
                    content = file.read().strip()  # Read and strip unnecessary whitespace
            except Exception as e:
                logger.error(f"No content found in {filename}. Due to error: {e}")
                content = ""
            sections_content.append((title, content))  # Append as tuple

        # Add the Executive Summary table to the presentation (second slide)
        self.create_presentation_functions.create_table_with_content(
            tid=self.tid, presentation=prs, sections_content=sections_content
        )

        # Save the presentation with both slides
        run_id = state["run_id"]
        sgr_output_pptx_path = os.path.join(SGR_OUTPUT_FOLDER, run_id)
        # if not os.path.exists(sgr_output_pptx_path):
        os.makedirs(sgr_output_pptx_path, exist_ok=True)
        prs.save(os.path.join(sgr_output_pptx_path, SGR_PPT_FILENAME))

        # Update trigger_flag_list to indicate successful completion of SGR
        # subgraph
        site_area = "SGR"
        trigger_flag_list = state["trigger_flag_list"]
        trigger_flag_list[site_area] = True
        return {
            "trigger_flag_list": trigger_flag_list,
            "output_files": {
                "combined_presentation": sgr_output_pptx_path,
            },
            "sgr_exec_agent_messages": AIMessage(
                name=f"{bold_start}CRM - generate_final_report node:{bold_end}",
                content=" Presentation created with " "PD summary and Executive summaries slides successfully!",
            ),
        }


class pdSummaryNodes:
    def __init__(self, pd_terms=pd_terms, PD_Term_DVTERM=PD_Term_DVTERM):
        self.pd_terms = pd_terms
        self.PD_Term_DVTERM = PD_Term_DVTERM
        self.sgr_pd_nodes_functions = sgrPDNodesFunctions()
        self.process_major_deviations_functions = self.sgr_pd_nodes_functions.processMajorDeviationsFunctions()
        self.generate_sgr_pd_findings_functions = self.sgr_pd_nodes_functions.generateSGRPDFindingsFunctions()

    def process_major_deviations(self, state: SGRSubGraphState):
        logger.debug("Calling function : process_major_deviations ")
        """Process and analyze major protocol deviations, saving results to CSV files."""
        # Load data
        pds = self.process_major_deviations_functions.load_data(state["input_filepath"]["pd_synthetic_data_filepath"])
        size = len(pds)
        # Part 1: Major Protocol Deviations
        major_devs = self.process_major_deviations_functions.filter_major_deviations(pds)
        grouped_counts = self.process_major_deviations_functions.group_and_count(
            major_devs, "Janssen Classification/DVDECOD", "Subject ID"
        )
        ordered_counts = self.process_major_deviations_functions.reindex_and_fill(grouped_counts, self.pd_terms)
        df_major_devs = self.process_major_deviations_functions.create_summary_table(
            ordered_counts,
            "Janssen Classification/DVDECOD",
            "Subject ID",
            size,
            "Subjects with at least 1 major protocol deviation",
        )

        self.process_major_deviations_functions.save_to_csv(
            df_major_devs,
            SGR_FILEPATHS["intermediate"]["dvdecode"],
            [
                "Major Protocol Deviation Analysis set: All Treated Through Week 48",
                f"Pooled Trial12/3 ({size} Subjects)",
            ],
        )
        logger.info("Generated intermediate file: dvdecode")

        top_subcat = ordered_counts.iloc[0]["Janssen Classification/DVDECOD"]
        top_subcat_action_taken = major_devs.loc[
            (major_devs["Janssen Classification/DVDECOD"] == top_subcat) & (major_devs["Action_Taken"].notnull()),
            ["Janssen Classification/DVDECOD", "Action_Taken", "Subject ID"],
        ]

        action_taken_counts = top_subcat_action_taken.groupby("Action_Taken")["Subject ID"].count().reset_index()
        action_taken_counts.columns = [
            f"Action_Taken: {top_subcat} ",
            "count_Subject_ID",
        ]

        action_taken_counts.to_csv(SGR_FILEPATHS["intermediate"]["dvdecode_AT"], index=False)
        logger.info("Generated intermediate file: dvdecode_AT")

        # Part 2: 'Other*' category deviations
        other_major_devs = self.process_major_deviations_functions.filter_major_deviations(
            pds[pds["Janssen Classification/DVDECOD"] == "Other*"]
        )
        grouped_other_counts = self.process_major_deviations_functions.group_and_count(
            other_major_devs, "PD Term/DVTERM", "Subject ID"
        )
        ordered_other_counts = self.process_major_deviations_functions.reindex_and_fill(grouped_other_counts, PD_Term_DVTERM)
        df_other_devs = self.process_major_deviations_functions.create_summary_table(
            ordered_other_counts,
            "PD Term/DVTERM",
            "Subject ID",
            size,
            "Subjects with at least 1 major protocol deviation in category of Other*",
        )
        self.process_major_deviations_functions.save_to_csv(
            df_other_devs,
            SGR_FILEPATHS["intermediate"]["pdterm"],
            [
                "Major Protocol Deviation Analysis set: All Treated",
                f"Pooled Trial12/3 ({size} Subjects)",
            ],
        )
        logger.info("Generated intermediate file: pdterm")

        other_top_subcat = ordered_other_counts.iloc[0]["PD Term/DVTERM"]
        other_top_subcat_action_taken = other_major_devs.loc[
            (other_major_devs["PD Term/DVTERM"] == other_top_subcat) & (other_major_devs["Action_Taken"].notnull()),
            ["PD Term/DVTERM", "Action_Taken", "Subject ID"],
        ]

        other_action_taken_counts = other_top_subcat_action_taken.groupby("Action_Taken")["Subject ID"].count().reset_index()
        other_action_taken_counts.columns = [
            f"Action_Taken: {other_top_subcat} ",
            "count_Subject ID",
        ]
        other_action_taken_counts.to_csv(SGR_FILEPATHS["intermediate"]["pdterm_AT"], index=False)
        logger.info("Generated intermediate file: pdterm_AT")
        state["intermediate_files"] = {
            "dvdecode": SGR_FILEPATHS["intermediate"]["dvdecode"],
            "dvdecode_AT": SGR_FILEPATHS["intermediate"]["dvdecode_AT"],
            "pdterm": SGR_FILEPATHS["intermediate"]["pdterm"],
            "pdterm_AT": SGR_FILEPATHS["intermediate"]["pdterm_AT"],
        }
        return {
            "intermediate_files": {
                "dvdecode": SGR_FILEPATHS["intermediate"]["dvdecode"],
                "dvdecode_AT": SGR_FILEPATHS["intermediate"]["dvdecode_AT"],
                "pdterm": SGR_FILEPATHS["intermediate"]["pdterm"],
                "pdterm_AT": SGR_FILEPATHS["intermediate"]["pdterm_AT"],
            },
            "sgr_agent_messages": AIMessage(
                name=f"{bold_start}CRM - process_major_deviations node:{bold_end}",
                content="Processing Major Deviations to"
                " create summaries at PD Term and DVDCODE levels ...."
                "\n  - Calculating sub-category (PD Term) wise distribution of subjects with atleast 1 major PD"
                "\n  - Calculating sub-category (DVDCODE) wise distribution of subjects with PD Term as 'Other'"
                "\n  - Documenting insights on actions taken for top deviation categories",
            ),
        }

    def generate_sgr_pd_findings(self, state: SGRSubGraphState):
        logger.debug("Calling function : generate_sgr_findings...")

        # Create a template for LangChain prompt
        prompt_template = sgr_pd_agent_prompts["GENERATE_SGR_FINDINGS_PROMPT"]

        # Initialize LangChain's PromptTemplate
        template = PromptTemplate(input_variables=["complete_prompt"], template=prompt_template)

        # Load the actual datasets
        dvdecode_data = read_file(state["intermediate_files"]["dvdecode"], "csv")
        pdterm_data = read_file(state["intermediate_files"]["pdterm"], "csv")
        dvdecode_AT_data = read_file(state["intermediate_files"]["dvdecode_AT"], "csv")
        pdterm_AT_data = read_file(state["intermediate_files"]["pdterm_AT"], "csv")

        # Prepare the prompt with data
        complete_prompt = self.generate_sgr_pd_findings_functions.prepare_prompt(
            "Analyze the two protocol deviation datasets.",
            dvdecode_data,
            pdterm_data,
            pdterm_AT_data,
            dvdecode_AT_data,
        )

        # Create the sequence using the new RunnableSequence (pipeline
        # approach)
        pipeline = template | llm

        # Run the pipeline
        response = pipeline.invoke({"complete_prompt": complete_prompt})

        output_text = response.content

        trial_id = TRIAL_ID
        sgr_pd_message = (
            "\nGenerating key"
            " observations for PD Term & DVDCODE summaries created earlier."
            "\n  - Processing and analyzing protocol deviation data using LLM for "
            "specific trial audit tasks."
            "\n  - Generating comprehensive reports summarizing major protocol deviations,"
            " including numerical summaries and qualitative insights."
            "consolidated outputs findings and focuses on top deviation categories "
            "and the actions taken to address these issues."
            f"\n  - Checking availability of activity findings for trial_id {trial_id} "
            "generated by PD Agent."
        )
        if "activity_findings" in state:
            if trial_id in state.get("activity_findings"):
                sgr_pd_message += "\n  - Activity findings found!"

                output_text += str(state["activity_findings"][trial_id][0])
        else:
            sgr_pd_message += ""

        with open(SGR_FILEPATHS["intermediate"]["keyobs"], "w") as file:
            for text in output_text:
                file.write(text)
        file.close()
        logger.info("Generated intermediate_file: keyobs")
        # reduce _AT table only for top category Action_taken dvdcode
        return {
            "sgr_agent_messages": AIMessage(
                name=f"{bold_start}CRM - generate_findings_agent:{bold_end}",
                content=sgr_pd_message,
            ),
            "llm_response": response.content,
            "intermediate_files": {"keyobs": SGR_FILEPATHS["intermediate"]["keyobs"]},
        }


class execSummaryNodes:
    def __init__(self, tid=trial_id):
        self.tid = tid

    def process_pd_section(self, state: SGRSubGraphState):
        logger.debug("Calling function : process_exec_summary_pd_section...")
        # Check if 'dvdecode' is in intermediate_files
        if "dvdecode" not in state["intermediate_files"]:
            raise KeyError("'dvdecode' is missing from intermediate_files")
        # Executive Summary -pd section
        df = read_file(state["intermediate_files"]["dvdecode"], "csv")

        response = f"{df.columns[1]}\n {df.iloc[0, 1]}"
        with open(SGR_FILEPATHS["intermediate"]["pd_section"], "w") as f:
            f.write(response)
        f.close()

        return {
            "intermediate_files": {"pd_section": SGR_FILEPATHS["intermediate"]["pd_section"]},
            "sgr_exec_agent_messages": AIMessage(
                name=f"{bold_start}CRM - process_pd_section node: {bold_end} ",
                content="Generated PD section findings ..."
                "\n  - Consolidating Major Deviations' Total Subjects and Number of Subjects with at least 1 major "
                "protocol deviation.",
            ),
        }

    def process_site_inspection_section(self, state: SGRSubGraphState):
        output = []
        insp_data = read_file(state["input_filepath"]["sponsor_inspection_data_filepath"], "csv")
        insp_data = insp_data.loc[insp_data["Protocol No"] == self.tid]

        quasar = read_file(state["input_filepath"]["quasar_data_filepath"], "csv")

        quasar = quasar.loc[quasar["trial_id"] == self.tid]

        si = set(insp_data["HA Country"].unique())
        output.append(f"{len(si)} site inspection ({', '.join(si)})")

        # can use LLM
        non_zero_critical = (quasar["critical_findings"] > 0).sum()

        # Count non-zero entries for 'major_findings'
        non_zero_major = (quasar["major_findings"] > 0).sum()

        output.append(f"{non_zero_critical} Critical and {non_zero_major} Major observations")

        sponsor_inspections = insp_data[insp_data["Inspection Type"] == "Sponsor GCP Inspection"].count().values[0]

        output.append(f"{sponsor_inspections} sponsor inspections")

        phrase = "Data integrity issue due to "

        # Filtering and counting titles that start with the specified phrase
        # can use LLM
        title_counts = quasar[quasar["finding_title"].str.startswith(phrase)]["finding_title"].value_counts()

        if title_counts.empty:
            output.append("No overall impact to data integrity of the submission")
        else:
            # Getting the count of the most common title starting with the
            # specified phrase
            most_common_count = title_counts.iloc[0]
            output.append(f"{most_common_count} {title_counts.index[0]} ")

        with open(
            SGR_FILEPATHS["intermediate"]["site_sponsor_inspections_section"],
            "w",
        ) as f:
            for text in output:
                f.write(text)
                f.write("\n")
        f.close()

        return {
            "intermediate_files": {
                "site_sponsor_inspections_section": SGR_FILEPATHS["intermediate"]["site_sponsor_inspections_section"]
            },
            "sgr_exec_agent_messages": AIMessage(
                name=f"{bold_start}CRM - process_site_inspection_section node: {bold_end} ",
                content="Generating findings"
                " from Sponsor Inspections data ..."
                "\n  - Processing and summarizing site inspection data, identifying the number and severity of "
                "observations across various locations."
                "\n  - Evaluating sponsor inspections and their implications on trial data integrity."
                "\n  - Compiling findings into a comprehensive report, documenting the impact and outcomes of site "
                "inspections.",
            ),
        }

    def process_significant_issue_section(self, state: SGRSubGraphState):
        # Site/Sponsor Related Self-Identified
        output = []
        trial_id = TRIAL_ID
        qis_sbs = read_file(state["input_filepath"]["qis_sbs_data_filepath"], "csv")
        site_count = (
            qis_sbs.loc[
                (qis_sbs["issue_self_identified"] == "Yes") & (qis_sbs["trial_id"] == self.tid),
                ["trial_id", "site_id", "issue_self_identified"],
            ]
            .groupby("site_id")["issue_self_identified"]
            .count()
            .index.nunique()
        )
        output.append(f"Site/Sponsor Related Self-Identified: {site_count}")

        serious_breach = read_file(state["input_filepath"]["serious_breach_data_filepath"], "csv")
        serious_breach_confirmed = (
            serious_breach.loc[
                serious_breach["trial_id"] == self.tid,
                ["trial_id", "serious_breach_confirmed"],
            ]
            .groupby(["serious_breach_confirmed"])
            .count()
        )

        if "Yes" in serious_breach_confirmed.index:
            value = serious_breach_confirmed.loc["Yes", "trial_id"]  # Get the count of 'Yes'
            output.append(f"{value} serious breach Notifications")

        else:
            output.append("No serious breach Notifications")

        # No overall impact to data integrity of the submission
        # can use LLM
        quasar = read_file(state["input_filepath"]["quasar_data_filepath"], "csv")
        quasar = quasar.loc[quasar["trial_id"] == self.tid]

        phrase = "Data integrity issue due to"

        # Filtering and counting titles that start with the specified phrase
        title_counts = quasar[quasar["finding_title"].str.startswith(phrase)]["finding_title"].value_counts()

        if title_counts.empty:
            output.append("No overall impact to data integrity of the submission")
        else:
            most_common_count = title_counts.iloc[0]
            output.append(f"{most_common_count} {title_counts.index[0]} ")

        sqi = read_file(state["input_filepath"]["sqi_data_filepath"], "csv")

        sqi = sqi.loc[sqi["trial_id"] == self.tid]
        # Add a column to indicate if it’s inspection-related or audit-related
        # can use LLM
        sqi["inspection_related"] = sqi["title"].apply(lambda title: 1 if "Site Inspection" in title else 0)
        sqi["audit_related"] = sqi["title"].apply(lambda title: 1 if "Audit Observation" in title else 0)

        significant_issues = ["Regulatory Compliance Issues"]
        closed_issues = sqi.loc[
            (sqi["trial_id"] == self.tid) & (sqi["record_state"] == "Closed") & sqi["issue_type"].isin(significant_issues)
        ]

        if closed_issues.empty:
            output.append("No significant issues resulted in site closures")
        else:
            output.append(f"{len(closed_issues)} significant issue(s) resulted in site closures ")

            inspection_related_count = closed_issues["inspection_related"].sum()
            audit_related_count = closed_issues["audit_related"].sum()
            output.append(f"Audit Related: {audit_related_count}\nInspection Related: {inspection_related_count}")

        # Filter rows where issue_type is a significant issue and record_state
        # is 'Closed'
        closed_significant_issues = sqi[(sqi["issue_type"].isin(significant_issues)) & (sqi["record_state"] == "Closed")]

        if closed_significant_issues.empty:
            output.append("No significant issues with a Closed record state found.")
        else:
            # Check how many of the closed significant issues have CAPA
            # completed
            completed_issues = closed_significant_issues[closed_significant_issues["CAPA_record_state"] == "completed"]

            # If all closed significant issues have CAPA completed
            if len(completed_issues) == len(closed_significant_issues):
                output.append("Significant issues having CAPAs completed.")
            else:
                completed_issues_count = len(completed_issues)
                closed_significant_issues_count = len(closed_significant_issues)
                output.append(
                    f"Only {completed_issues_count} out of {closed_significant_issues_count} significant issues"
                    " have CAPAs completed."
                )

        with open(
            SGR_FILEPATHS["intermediate"]["significant_issue_escalations_section"],
            "w",
        ) as f:
            for text in output:
                f.write(text)
                f.write("\n")
        f.close()

        return {
            "intermediate_files": {
                "significant_issue_escalations_section": SGR_FILEPATHS["intermediate"][
                    "significant_issue_escalations_section"
                ]
            },
            "sgr_exec_agent_messages": AIMessage(
                name=f"{bold_start}CRM - process_significant_issue_escalation_section node: {bold_end} ",
                content="Generating findings"
                " from Significant Issue Escalations data ..."
                f"\n  - Processing and summarizing significant issue escalation data for trial_id {trial_id} ."
                " \n  - Outputs "
                "include assessments of serious breach notifications, site/sponsor self-identified issues, data "
                "integrity impacts, and audit/inspection-related findings are consolidated.",
            ),
        }

    def process_qa_audit_section(self, state: SGRSubGraphState):
        output = []

        # 0 Critical Audit observations
        # Can use LLM
        sqi = read_file(SGR_FILEPATHS["input"]["sqi_data"], "csv")

        sqi = sqi[sqi["trial_id"] == self.tid]
        critical_observations = sqi["title"].apply(lambda title: 1 if "Critical GCP Audit Observation" in title else 0)

        output.append(f"{critical_observations.sum()} Critical Audit observations")

        # Clinical Investigators: 4 Audits conducted (2 routine, 0 focuse)
        qaa = read_file(SGR_FILEPATHS["input"]["quasar_data"], "csv")
        qaa = qaa.loc[qaa["trial_id"] == self.tid]

        # Filter for Clinical Investigators
        clinical_investigators_df = qaa[qaa["audit_type"] == "Clinical Investigators"]

        # Get the total number of Clinical Investigators audits
        total_audits = len(clinical_investigators_df)

        # Count how many are Routine and how many are Focuse
        audit_reason_counts = clinical_investigators_df["audit_reason"].value_counts()

        # Get the counts for Routine and Focuse (defaulting to 0 if not
        # present)
        routine_count = audit_reason_counts.get("Routine", 0)
        focuse_count = audit_reason_counts.get("focuse", 0)

        output.append(
            f"Clinical Investigators: {total_audits} Audits conducted ({routine_count} routine, {focuse_count} focuse)"
        )
        # # 3 System Audits: 3 Major
        # Group audits by type and count them
        audits = qaa.groupby(["audit_type"])["audit_id"].count().reset_index()

        # Check if 'System' audits exist
        if not audits[audits["audit_type"] == "System"].empty:
            sa = audits[audits["audit_type"] == "System"].values[0]
            audit_output = f"{sa[1]} {sa[0]} Audits"
        else:
            audit_output = ""

        # Group by finding classification and count
        classification = qaa[qaa["audit_type"] == "System"].groupby(["finding_classification"]).count().reset_index()

        # Check if 'Major' classification exists within 'System' audits
        if not classification[classification["finding_classification"] == "Major"].empty:
            mj = classification[classification["finding_classification"] == "Major"].values[0]
            # Append details about 'Major' findings only if they exist
            audit_output += f": {mj[1]} {mj[0]}"

        # Append to output list
        output.append(audit_output)

        # No Significant imapct on patient safety or data integrity
        # Need to use LLM

        with open(SGR_FILEPATHS["intermediate"]["qa_audit_section"], "w") as f:
            for text in output:
                f.write(text)
                f.write("\n")
        f.close()
        trial_id = TRIAL_ID
        return {
            "intermediate_files": {"qa_audit_section": SGR_FILEPATHS["intermediate"]["qa_audit_section"]},
            "sgr_exec_agent_messages": AIMessage(
                name=f"{bold_start}CRM - process_qa_audit_section node: {bold_end} ",
                content="Generated findings"
                " from QA Audit data ..."
                f"\n  - Processed and summarized QA audit data for trial_id {trial_id} ; findings including "
                "critical observations, audit types, and impacts on patient safety and data integrity"
                " are consolidated.",
            ),
        }
