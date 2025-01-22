import pandas as pd
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from ....utils.helpers import read_file
from ....utils.log_setup import get_logger

# Get the same logger instance set up earlier
logger = get_logger()


class sgrPDNodesFunctions:
    def __init__(self):
        logger.debug("Initialising sgrPDNodesFunctions ...")
        pass

    class processMajorDeviationsFunctions:
        def __init__(self):
            logger.debug("Initialising processMajorDeviationsFunctions ...")
            pass

        # Define functions for repetitive operations
        def load_data(
            self,
            file_path: str,
        ):
            """Load the CSV file and return the DataFrame."""
            df = read_file(file_path, "csv")
            return df

        def filter_major_deviations(self, df, category_col="Confirmed Category", category_value="MAJOR"):
            """Filter the DataFrame for rows with major deviations."""
            return df[df[category_col] == category_value]

        def group_and_count(self, df, group_col, count_col):
            """Group by a specific column and count occurrences in another column."""
            return df.groupby(group_col)[count_col].count()

        def reindex_and_fill(self, grouped_counts, terms, fill_value=1):
            """Reindex the grouped counts DataFrame based on predefined terms, filling missing values."""
            return grouped_counts.reindex(terms, fill_value=fill_value).reset_index()

        def create_summary_table(self, ordered_counts, terms_col, count_col, total_subjects, label):
            """Create a summary DataFrame with count and percentage format."""
            df = pd.DataFrame(
                {
                    terms_col: ordered_counts[terms_col],
                    "Subjects with at least 1 major protocol deviation": (
                        ordered_counts[count_col].astype(str)
                        + " ("
                        + (ordered_counts[count_col] / total_subjects * 100).round(1).astype(str)
                        + "%)"
                    ),
                }
            )

            # Insert summary row
            total_deviations = ordered_counts[count_col].sum()
            summary_row = [
                label,
                f"{total_deviations} ({(total_deviations / total_subjects * 100).round(1)}%)",
            ]
            df.loc[-1] = summary_row
            df.index = df.index + 1  # Adjust index
            df = df.sort_index()  # Sort index to ensure proper order
            return df

        def save_to_csv(self, df, file_name, columns):
            """Save the DataFrame to a CSV file."""
            df.columns = columns
            df.to_csv(file_name, index=False)

    class generateSGRPDFindingsFunctions:
        def __init__(self):
            pass

        def prepare_prompt(self, prompt, df1, df2, df3, df4):
            # Function to prepare prompt with DataFrame data
            # Convert DataFrames to a string format
            df1_string = df1.to_string()
            df2_string = df2.to_string()
            df3_string = df3.to_string()
            df4_string = df4.to_string()

            # Combine prompt and data
            complete_prompt = (
                f"{prompt}\n\nFirst dataset:\n{df1_string}\n\nSecond dataset:\n{df2_string}\n"
                f"\nThird dataset:\n{df3_string}\n\nFourth dataset:\n{df4_string}"
            )
            return complete_prompt

    class createSGRPDSlideFunctions:
        def __init__(self):
            pass

        def add_table_to_slide(
            slide,
            df,
            left,
            top,
            width,
            height,
            header_color=RGBColor(0, 0, 139),
            row_color=RGBColor(173, 216, 230),
        ):
            """Add a DataFrame as a table to the slide at the given position with specified colors."""

            # Create table with rows and columns based on DataFrame size
            rows, cols = df.shape
            table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

            # Set column widths
            for col in range(cols):
                table.columns[col].width = width // cols

            # Set header row
            for col, column_name in enumerate(df.columns):
                cell = table.cell(0, col)
                cell.text = column_name
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.size = Pt(12)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text for dark blue header

            # Add DataFrame rows to the table
            for row in range(rows):
                for col in range(cols):
                    cell = table.cell(row + 1, col)
                    cell.text = str(df.iloc[row, col])
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_color
                    cell.text_frame.paragraphs[0].font.size = Pt(10)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            return table


class createSGRPresentationFunctions:
    def __init__(self):
        pass

    # Function to set the font properties (size, color, bold)
    def set_font(self, cell, font_size, font_color, bold=False):
        text_frame = cell.text_frame
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
                run.font.color.rgb = font_color
                run.font.bold = bold

    # Function to create a 2-column table with content for Slide 2
    def create_table_with_content(self, tid, presentation, sections_content):
        slide_layout = presentation.slide_layouts[5]  # Blank slide
        slide = presentation.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = "Executive Summary - Trials 12/3"
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.5)
        title_shape.width = Inches(4)
        title_shape.height = Inches(0.5)
        title_shape.text_frame.paragraphs[0].font.size = Pt(15)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 139)  # Dark blue color

        # Define table dimensions
        rows = len(sections_content) + 1  # Including header row
        cols = 2  # 2 columns (Section and Details)

        # Set table position and size
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9.5)
        height = Inches(4.5)

        # Add table
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Define colors
        header_color = RGBColor(0, 0, 139)  # Dark Blue for header
        row_color = RGBColor(173, 216, 230)  # Light Blue for rows
        first_column_color = RGBColor(0, 0, 139)  # Dark Blue for the first column
        white_text_color = RGBColor(255, 255, 255)  # White text

        # Set column widths
        table.columns[0].width = Inches(1)  # First column (title) width
        table.columns[1].width = Inches(7.0)  # Second column (description) width

        # Set header row content
        table.cell(0, 0).text = "Section"
        table.cell(0, 1).text = "Findings "

        # Apply header color
        for col in range(2):
            cell = table.cell(0, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            self.set_font(cell, font_size=12, font_color=white_text_color, bold=True)

        # Fill table rows with content and apply colors
        for row_idx, (section, details) in enumerate(sections_content, start=1):
            # First column (section)
            table.cell(row_idx, 0).text = section
            table.cell(row_idx, 0).fill.solid()
            table.cell(row_idx, 0).fill.fore_color.rgb = first_column_color
            self.set_font(
                table.cell(row_idx, 0),
                font_size=12,
                font_color=white_text_color,
                bold=True,
            )

            # Second column (details)
            table.cell(row_idx, 1).text = details
            table.cell(row_idx, 1).fill.solid()
            table.cell(row_idx, 1).fill.fore_color.rgb = row_color
            self.set_font(
                table.cell(row_idx, 1),
                font_size=12,
                font_color=RGBColor(0, 0, 0),
            )  # Black text

    # Function to add a table to a PowerPoint slide for Slide 1
    def add_table_to_slide(
        self,
        slide,
        df,
        left,
        top,
        width,
        height,
        header_color=RGBColor(0, 0, 139),
        row_color=RGBColor(173, 216, 230),
    ):
        """Add a DataFrame as a table to the slide at the given position with specified colors."""

        # Create table with rows and columns based on DataFrame size
        rows, cols = df.shape
        table = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

        # Set column widths
        for col in range(cols):
            table.columns[col].width = width // cols

        # Set header row
        for col, column_name in enumerate(df.columns):
            cell = table.cell(0, col)
            cell.text = column_name
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text for dark blue header

        # Add DataFrame rows to the table
        for row in range(rows):
            for col in range(cols):
                cell = table.cell(row + 1, col)
                cell.text = str(df.iloc[row, col])
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_color
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        return table
